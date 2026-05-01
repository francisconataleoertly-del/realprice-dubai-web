from pathlib import Path
from shutil import copy2

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(r"C:\Users\franc\realprice-dubai-web")
DESKTOP = Path(r"C:\Users\franc\OneDrive\Desktop")
DOWNLOADS = Path(r"C:\Users\franc\Downloads")
BRAND = REPO / "public" / "brand"
DUBAI = REPO / "public" / "dubai-slides"

OUT_DESKTOP = DESKTOP / "FonatProp_Broker_Demo_5SLIDES_VISUAL.pptx"
OUT_DOWNLOADS = DOWNLOADS / "FonatProp_Broker_Demo_5SLIDES_VISUAL.pptx"

DEMO_URL = "https://fonatprop.com/broker-demo"
WIDGET_URL = "https://fonatprop.com/widget/embed.js"
LOGO = BRAND / "fonatprop-logo-lockup-transparent.png"
QR_PATH = BRAND / "broker-demo-qr.png"

qr = qrcode.QRCode(border=1, box_size=12)
qr.add_data(DEMO_URL)
qr.make(fit=True)
qr.make_image(fill_color="white", back_color="#050914").convert("RGB").save(QR_PATH)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

WHITE = RGBColor(246, 248, 252)
MUTED = RGBColor(177, 186, 201)
SOFT = RGBColor(118, 132, 154)
BLUE = RGBColor(59, 130, 246)
DARK = RGBColor(3, 7, 14)
PANEL = RGBColor(8, 14, 28)


def add_text(slide, text, left, top, width, height, size=20, color=WHITE, bold=False, font="Aptos", italic=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    return box


def add_background(slide):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = DARK
    rect.line.fill.background()


def add_logo(slide):
    slide.shapes.add_picture(str(LOGO), Inches(0.64), Inches(0.42), width=Inches(2.18))


def add_photo(slide, filename, caption):
    # Big, visible Dubai image panel. No dark overlay on top of the photo.
    left, top, width, height = Inches(7.25), Inches(0.70), Inches(5.38), Inches(5.75)
    slide.shapes.add_picture(str(DUBAI / filename), left, top, width=width, height=height)
    border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    border.fill.background()
    border.line.color.rgb = RGBColor(58, 76, 104)
    border.line.transparency = 35
    cap = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + height - Inches(0.54), width, Inches(0.54))
    cap.fill.solid()
    cap.fill.fore_color.rgb = DARK
    cap.fill.transparency = 18
    cap.line.fill.background()
    add_text(slide, caption.upper(), left + Inches(0.22), top + height - Inches(0.40), width - Inches(0.44), Inches(0.22), 8, SOFT, False, "Consolas")


def label(slide, text):
    add_text(slide, text.upper(), Inches(0.70), Inches(1.45), Inches(4.7), Inches(0.28), 9, SOFT, False, "Consolas")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.70), Inches(1.82), Inches(0.75), Inches(0.025))
    line.fill.solid()
    line.fill.fore_color.rgb = BLUE
    line.line.fill.background()


def title(slide, line1, line2):
    add_text(slide, line1, Inches(0.70), Inches(2.05), Inches(6.0), Inches(0.70), 38, WHITE, False, "Georgia")
    add_text(slide, line2, Inches(0.70), Inches(2.70), Inches(6.0), Inches(0.78), 38, RGBColor(168, 171, 180), False, "Georgia", True)


def body(slide, text):
    add_text(slide, text, Inches(0.76), Inches(3.78), Inches(5.76), Inches(1.06), 16, MUTED, False, "Aptos")


def bullets(slide, items, top=5.22):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.78), Inches(y + 0.08), Inches(0.09), Inches(0.09))
        dot.fill.solid()
        dot.fill.fore_color.rgb = BLUE
        dot.line.fill.background()
        add_text(slide, item, Inches(1.00), Inches(y), Inches(5.35), Inches(0.38), 13.5, MUTED, False, "Aptos")
        y += 0.47


def footer(slide, num):
    add_text(slide, f"{num:02d} / FonatProp", Inches(11.56), Inches(6.96), Inches(1.1), Inches(0.20), 7, SOFT, False, "Consolas", align=PP_ALIGN.RIGHT)


def add_standard_slide(num, photo, photo_caption, section, headline1, headline2, copy, bullet_items):
    slide = prs.slides.add_slide(BLANK)
    add_background(slide)
    add_logo(slide)
    add_photo(slide, photo, photo_caption)
    label(slide, f"{num:02d} / {section}")
    title(slide, headline1, headline2)
    body(slide, copy)
    bullets(slide, bullet_items)
    footer(slide, num)


add_standard_slide(
    1,
    "05-downtown-night.jpg",
    "Downtown Dubai",
    "Brokerage growth demo",
    "Two tools.",
    "One commercial goal.",
    "FonatProp gives brokerages a private valuation tool for agents, plus a public website widget designed to generate qualified property-owner inquiries.",
    ["Private exact valuation for the brokerage", "Public widget for marketing and contact capture", "The agent controls the precise valuation conversation"],
)

add_standard_slide(
    2,
    "01-marina-skyline.jpg",
    "Dubai Marina",
    "The problem",
    "Agency websites",
    "need stronger conversion.",
    "Most brokerage websites look good, but they do not create enough reasons for property owners, buyers or investors to start a conversation.",
    ["Traffic leaves without contacting an agent", "Owners need a reason to share property details", "The website should create qualified inquiries"],
)

add_standard_slide(
    3,
    "09-palm-aerial.jpg",
    "Palm Jumeirah",
    "Private valuation",
    "Exact numbers",
    "stay with the brokerage.",
    "The professional valuation surface is private. Only the brokerage team gets access to the detailed AI estimate, comparables and confidence range.",
    ["Controlled access for agents and managers", "Real Dubai transaction evidence", "Designed for pricing, meetings and follow-up"],
)

add_standard_slide(
    4,
    "business-bay.jpg",
    "Business Bay",
    "Public widget",
    "A marketing hook",
    "not the full valuation.",
    "The website widget is a conversion tool. It gives visitors a broad, non-final indication and routes the inquiry to the brokerage for the real valuation call.",
    ["Captures name, phone, email and property details", "Shows only a general range or market benchmark", "Sends the inquiry to WhatsApp, email, Sheets or CRM"],
)

slide = prs.slides.add_slide(BLANK)
add_background(slide)
add_logo(slide)
add_photo(slide, "02-burj-khalifa.jpg", "Live demo")
label(slide, "05 / Next step")
title(slide, "Show the flow.", "Then start the pilot.")
body(slide, "The demo separates the two products clearly: private AI valuation for the brokerage, and a public widget that captures client inquiries.")

panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.76), Inches(5.06), Inches(5.95), Inches(1.35))
panel.fill.solid()
panel.fill.fore_color.rgb = PANEL
panel.line.color.rgb = RGBColor(70, 88, 120)
panel.line.transparency = 30
add_text(slide, "LIVE BROKER DEMO", Inches(1.08), Inches(5.30), Inches(2.7), Inches(0.23), 8, SOFT, False, "Consolas")
add_text(slide, DEMO_URL, Inches(1.08), Inches(5.66), Inches(4.75), Inches(0.30), 15, WHITE, True, "Aptos")
add_text(slide, f"Widget script: {WIDGET_URL}", Inches(1.08), Inches(6.08), Inches(5.20), Inches(0.22), 8.5, RGBColor(186, 211, 255), False, "Consolas")
btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.92), Inches(5.46), Inches(1.35), Inches(0.50))
btn.fill.solid()
btn.fill.fore_color.rgb = RGBColor(255, 255, 255)
btn.line.color.rgb = RGBColor(255, 255, 255)
btn.click_action.hyperlink.address = DEMO_URL
tf = btn.text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "OPEN"
r.font.size = Pt(10.5)
r.font.bold = True
r.font.color.rgb = DARK
if QR_PATH.exists():
    slide.shapes.add_picture(str(QR_PATH), Inches(6.02), Inches(5.24), width=Inches(0.86))
footer(slide, 5)

prs.save(str(OUT_DOWNLOADS))
copy2(OUT_DOWNLOADS, OUT_DESKTOP)
print(OUT_DOWNLOADS)
print(OUT_DESKTOP)
print(len(prs.slides))
print(round(OUT_DOWNLOADS.stat().st_size / 1024 / 1024, 2))
