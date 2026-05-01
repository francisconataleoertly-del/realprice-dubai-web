from pathlib import Path
from shutil import copy2

import qrcode
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

REPO = Path(r"C:\Users\franc\realprice-dubai-web")
DESKTOP = Path(r"C:\Users\franc\OneDrive\Desktop")
DOWNLOADS = Path(r"C:\Users\franc\Downloads")
BRAND = REPO / "public" / "brand"
SLIDES_DIR = REPO / "public" / "dubai-slides"

OUT_DOWNLOADS = DOWNLOADS / "FonatProp_Dubai_Brokerages_CLEAN_6SLIDES.pptx"
OUT_DESKTOP = DESKTOP / "FonatProp_CLEAN_6SLIDES.pptx"
DEMO_URL = "https://fonatprop.com/broker-demo"
WIDGET_URL = "https://fonatprop.com/widget/embed.js"
LOGO = BRAND / "fonatprop-logo-lockup-transparent.png"
QR_PATH = BRAND / "broker-demo-qr.png"

qr = qrcode.QRCode(border=1, box_size=12)
qr.add_data(DEMO_URL)
qr.make(fit=True)
qr.make_image(fill_color="white", back_color="#050914").convert("RGB").save(QR_PATH)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = prs.slide_width, prs.slide_height

COLORS = {
    "white": RGBColor(246, 248, 252),
    "muted": RGBColor(174, 183, 198),
    "soft": RGBColor(125, 138, 158),
    "blue": RGBColor(59, 130, 246),
    "dark": RGBColor(3, 7, 14),
    "panel": RGBColor(7, 12, 24),
}

BACKGROUNDS = [
    "05-downtown-night.jpg",
    "01-marina-skyline.jpg",
    "09-palm-aerial.jpg",
    "business-bay.jpg",
    "03-burj-al-arab.jpg",
    "02-burj-khalifa.jpg",
]


def add_bg(slide, image_name: str, dark: int = 36) -> None:
    slide.shapes.add_picture(str(SLIDES_DIR / image_name), 0, 0, width=W, height=H)
    overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = COLORS["dark"]
    overlay.fill.transparency = dark
    overlay.line.fill.background()
    side = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(7.8), H)
    side.fill.solid()
    side.fill.fore_color.rgb = COLORS["dark"]
    side.fill.transparency = 18
    side.line.fill.background()


def textbox(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    size: float = 24,
    bold: bool = False,
    color=None,
    font: str = "Aptos",
    align=PP_ALIGN.LEFT,
    italic: bool = False,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or COLORS["white"]
    return box


def logo_header(slide) -> None:
    slide.shapes.add_picture(str(LOGO), Inches(0.70), Inches(0.42), width=Inches(2.18))


def section_label(slide, text: str, top=1.46) -> None:
    textbox(slide, text.upper(), Inches(0.78), Inches(top), Inches(5.7), Inches(0.32), 9, False, COLORS["soft"], "Consolas")
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.78), Inches(top + 0.34), Inches(0.76), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["blue"]
    line.line.fill.background()


def title(slide, main: str, accent: str, top=1.92) -> None:
    textbox(slide, main, Inches(0.72), Inches(top), Inches(7.4), Inches(0.86), 44, False, COLORS["white"], "Georgia")
    textbox(slide, accent, Inches(0.72), Inches(top + 0.72), Inches(7.4), Inches(0.9), 44, False, RGBColor(170, 172, 178), "Georgia", italic=True)


def body(slide, text: str, top, width=5.9) -> None:
    textbox(slide, text, Inches(0.78), top, Inches(width), Inches(1.1), 17, False, COLORS["muted"], "Aptos")


def metric(slide, value: str, label: str, left, top) -> None:
    textbox(slide, value, left, top, Inches(1.8), Inches(0.42), 25, True, COLORS["white"], "Georgia")
    textbox(slide, label.upper(), left, top + Inches(0.42), Inches(2.2), Inches(0.28), 7.8, False, COLORS["soft"], "Consolas")


def panel(slide, left, top, width, height):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = COLORS["panel"]
    shp.fill.transparency = 7
    shp.line.color.rgb = RGBColor(70, 88, 120)
    shp.line.transparency = 40
    return shp


def bullet(slide, text: str, left, top, width=4.4) -> None:
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top + Inches(0.08), Inches(0.09), Inches(0.09))
    dot.fill.solid()
    dot.fill.fore_color.rgb = COLORS["blue"]
    dot.line.fill.background()
    textbox(slide, text, left + Inches(0.22), top, Inches(width), Inches(0.52), 14, False, COLORS["muted"], "Aptos")


def add_content_slide(num: int, bg: str, label: str, main: str, accent: str, copy: str, bullets=None, metrics=None):
    slide = prs.slides.add_slide(BLANK)
    add_bg(slide, bg)
    logo_header(slide)
    section_label(slide, f"{num:02d} / {label}")
    title(slide, main, accent)
    body(slide, copy, Inches(3.80), width=6.3)
    if bullets:
        panel(slide, Inches(8.25), Inches(1.55), Inches(4.15), Inches(3.65))
        textbox(slide, "Broker value", Inches(8.55), Inches(1.86), Inches(3.3), Inches(0.35), 10, False, COLORS["soft"], "Consolas")
        y = 2.38
        for item in bullets:
            bullet(slide, item, Inches(8.55), Inches(y), width=3.35)
            y += 0.72
    if metrics:
        x = 0.78
        for value, label_text in metrics:
            metric(slide, value, label_text, Inches(x), Inches(6.20))
            x += 2.1
    textbox(slide, f"{num:02d} / FonatProp", Inches(11.55), Inches(6.95), Inches(1.1), Inches(0.22), 7, False, COLORS["soft"], "Consolas", PP_ALIGN.RIGHT)


add_content_slide(
    1,
    BACKGROUNDS[0],
    "Brokerage revenue demo",
    "Dubai brokerages",
    "need seller leads.",
    "FonatProp adds AI-powered valuation and seller lead capture to a brokerage website, without rebuilding the site.",
    ["Instant value proposition for property owners", "One embedded widget for the broker website", "Agent keeps the exact valuation conversation"],
    [("234K+", "verified Dubai transactions"), ("1 script", "website install"), ("Live", "demo ready")],
)

add_content_slide(
    2,
    BACKGROUNDS[1],
    "Problem",
    "Most agency websites",
    "do not convert owners.",
    "Visitors browse listings, but owners who may sell or rent often leave without becoming a lead. The broker loses the highest-value opportunity: the property owner.",
    ["Traffic without owner data", "No valuation hook on the website", "Manual follow-up starts too late", "Competitors look more digital"],
)

add_content_slide(
    3,
    BACKGROUNDS[2],
    "Valuation engine",
    "A smarter first answer",
    "for Dubai property owners.",
    "The valuation tool uses real Dubai transaction evidence to estimate a property range by address, area, bedrooms and building context.",
    ["Address and building-aware estimates", "Dubai zones and transaction anchors", "Confidence range instead of false certainty", "Designed to support the broker, not replace them"],
)

add_content_slide(
    4,
    BACKGROUNDS[3],
    "Website widget",
    "Turn the website",
    "into a seller lead machine.",
    "The widget captures name, email, phone and property details, then sends the enriched lead directly to WhatsApp, email, Zapier, Make, Google Sheets or CRM.",
    ["Lead capture before the estimate", "General range for the visitor", "Agent handoff after intent is confirmed", "English and Arabic-ready experience"],
)

add_content_slide(
    5,
    BACKGROUNDS[4],
    "Pilot proposal",
    "A simple pilot",
    "for one brokerage website.",
    "Start with one brokerage, one website and one lead workflow. Prove that valuation creates conversations with sellers, then scale to more agents.",
    ["Install widget on homepage or seller page", "Route leads to agent WhatsApp/email", "Weekly lead report", "Optional CRM integration after validation"],
)

slide = prs.slides.add_slide(BLANK)
add_bg(slide, BACKGROUNDS[5], dark=30)
logo_header(slide)
section_label(slide, "06 / Live demo", top=1.28)
title(slide, "See it live.", "Then install it.", top=1.72)
body(slide, "For the meeting, open the broker demo first. It shows the internal valuation surface and the embeddable website widget in one focused sales flow.", Inches(3.52), width=6.3)
panel(slide, Inches(7.20), Inches(1.42), Inches(4.92), Inches(4.80))
textbox(slide, "LIVE BROKER DEMO", Inches(7.62), Inches(1.86), Inches(3.7), Inches(0.34), 10, False, COLORS["soft"], "Consolas")
textbox(slide, DEMO_URL, Inches(7.62), Inches(2.34), Inches(4.0), Inches(0.42), 18, True, COLORS["white"], "Aptos")
textbox(slide, "WIDGET SCRIPT", Inches(7.62), Inches(3.10), Inches(3.7), Inches(0.34), 10, False, COLORS["soft"], "Consolas")
textbox(slide, WIDGET_URL, Inches(7.62), Inches(3.56), Inches(4.0), Inches(0.42), 13, True, RGBColor(191, 217, 255), "Consolas")
if QR_PATH.exists():
    qr_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.86), Inches(4.32), Inches(1.62), Inches(1.62))
    qr_bg.fill.solid()
    qr_bg.fill.fore_color.rgb = RGBColor(5, 9, 18)
    qr_bg.line.color.rgb = RGBColor(255, 255, 255)
    qr_bg.line.transparency = 45
    slide.shapes.add_picture(str(QR_PATH), Inches(9.98), Inches(4.44), width=Inches(1.38))
btn = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.62), Inches(4.42), Inches(1.85), Inches(0.56))
btn.fill.solid()
btn.fill.fore_color.rgb = RGBColor(255, 255, 255)
btn.line.color.rgb = RGBColor(255, 255, 255)
btn.click_action.hyperlink.address = DEMO_URL
tf = btn.text_frame
tf.clear()
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "OPEN DEMO"
r.font.size = Pt(11)
r.font.bold = True
r.font.color.rgb = COLORS["dark"]
textbox(slide, "06 / FonatProp", Inches(11.55), Inches(6.95), Inches(1.1), Inches(0.22), 7, False, COLORS["soft"], "Consolas", PP_ALIGN.RIGHT)

prs.save(str(OUT_DOWNLOADS))
copy2(OUT_DOWNLOADS, OUT_DESKTOP)
print(OUT_DOWNLOADS)
print(OUT_DESKTOP)
print(len(prs.slides))
print(round(OUT_DOWNLOADS.stat().st_size / 1024 / 1024, 2))
