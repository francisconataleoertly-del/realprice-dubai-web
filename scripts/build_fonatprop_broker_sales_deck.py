from __future__ import annotations

import shutil
import tempfile
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT_DESKTOP = Path.home() / "OneDrive" / "Desktop" / "FonatProp_Broker_Deck_PREMIUM.pptx"
OUT_DOWNLOADS = Path.home() / "Downloads" / "FonatProp_Broker_Deck_PREMIUM.pptx"
OUT_DATA_LAKE = Path.home() / "FonatProp_Data_Lake" / "05_decks" / "FonatProp_Broker_Deck_PREMIUM.pptx"
DEMO_URL = "https://fonatprop.com/broker-demo"
CONTACT_PHONE_DISPLAY = "+54 9 11 2640-9578"
CONTACT_WHATSAPP_URL = "https://wa.me/5491126409578"
CONTACT_EMAIL = "fonatprop@gmail.com"
CONTACT_EMAIL_URL = f"mailto:{CONTACT_EMAIL}"
DECK_MEDIA = ROOT / "public" / "deck-media"

W, H = 1920, 1080
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG_IMAGES = [
    ROOT / "public" / "dubai-slides" / "02-burj-khalifa.jpg",
    ROOT / "public" / "dubai-slides" / "01-marina-skyline.jpg",
    ROOT / "public" / "dubai-slides" / "04-marina-night.jpg",
    ROOT / "public" / "dubai-slides" / "07-marina-aerial.jpg",
    ROOT / "public" / "dubai-slides" / "03-burj-al-arab.jpg",
    ROOT / "public" / "dubai-slides" / "05-downtown-night.jpg",
    ROOT / "public" / "dubai-slides" / "06-sunset-silhouette.jpg",
]

LOGO = ROOT / "public" / "brand" / "fonatprop-logo-lockup-transparent.webp"
MARK = ROOT / "public" / "brand" / "fonatprop-final-icon.png"
QR = ROOT / "public" / "brand" / "broker-demo-qr.png"
FOUNDER_PHOTO = Path(
    r"C:\Users\franc\Downloads\WhatsApp Image 2026-05-03 at 20.41.22.jpeg"
)
VALUATION_VIDEO = Path(
    r"C:\Users\franc\Videos\Captures\FonatProp Broker Demo _ Valuation + Widget - Google Chrome 2026-04-29 20-28-19.mp4"
)
WIDGET_VIDEO = Path(
    r"C:\Users\franc\Videos\Captures\FonatProp Broker Demo _ Valuation + Widget - Google Chrome 2026-04-29 20-29-25.mp4"
)
VALUATION_POSTER = DECK_MEDIA / "valuation-poster.png"
WIDGET_POSTER = DECK_MEDIA / "widget-poster.png"

WHITE = (247, 249, 255)
SOFT = (201, 208, 226)
MUTED = (139, 150, 176)
BLUE = (92, 132, 255)
CYAN = (83, 223, 255)
VIOLET = (143, 92, 255)
MINT = (88, 247, 205)
GOLD = (235, 196, 105)
CARD = (7, 10, 23, 222)
CARD2 = (10, 14, 30, 238)
INK = (3, 5, 10)


def font(name: str, size: int) -> ImageFont.FreeTypeFont:
    candidates = {
        "serif": [
            "C:/Windows/Fonts/bahnschrift.ttf",
            "C:/Windows/Fonts/segoeuib.ttf",
            "C:/Windows/Fonts/arialbd.ttf",
        ],
        "serif_i": [
            "C:/Windows/Fonts/segoeuisl.ttf",
            "C:/Windows/Fonts/segoeuii.ttf",
            "C:/Windows/Fonts/segoeuil.ttf",
        ],
        "sans": [
            "C:/Windows/Fonts/segoeui.ttf",
            "C:/Windows/Fonts/aptos.ttf",
            "C:/Windows/Fonts/arial.ttf",
        ],
        "sans_b": [
            "C:/Windows/Fonts/bahnschrift.ttf",
            "C:/Windows/Fonts/segoeuib.ttf",
            "C:/Windows/Fonts/arialbd.ttf",
        ],
        "mono": [
            "C:/Windows/Fonts/consolab.ttf",
            "C:/Windows/Fonts/consola.ttf",
        ],
    }
    for path in candidates[name]:
        if Path(path).exists():
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


SERIF_88 = font("serif", 88)
SERIF_78 = font("serif", 78)
SERIF_66 = font("serif", 66)
SERIF_56 = font("serif", 56)
SERIF_I_76 = font("serif_i", 76)
SERIF_I_58 = font("serif_i", 58)
SANS_32 = font("sans", 32)
SANS_28 = font("sans", 28)
SANS_24 = font("sans", 24)
SANS_22 = font("sans", 22)
SANS_20 = font("sans", 20)
SANS_18 = font("sans", 18)
SANS_64 = font("sans", 64)
SANS_52 = font("sans", 52)
SANS_B_34 = font("sans_b", 34)
SANS_B_28 = font("sans_b", 28)
SANS_B_22 = font("sans_b", 22)
SANS_B_18 = font("sans_b", 18)
SANS_B_74 = font("sans_b", 74)
SANS_B_56 = font("sans_b", 56)
MONO_18 = font("mono", 18)
MONO_15 = font("mono", 15)


def cover_image(path: Path) -> Image.Image:
    img = Image.open(path).convert("RGB")
    ratio = max(W / img.width, H / img.height)
    new_size = (int(img.width * ratio), int(img.height * ratio))
    img = img.resize(new_size, Image.Resampling.LANCZOS)
    left = (img.width - W) // 2
    top = (img.height - H) // 2
    return img.crop((left, top, left + W, top + H))


def premium_bg(path: Path, dark: int = 88) -> Image.Image:
    img = cover_image(path)
    img = ImageEnhance.Brightness(img).enhance(0.96)
    img = ImageEnhance.Contrast(img).enhance(1.14)
    overlay = Image.new("RGBA", (W, H), (3, 5, 12, max(56, dark - 14)))
    img = Image.alpha_composite(img.convert("RGBA"), overlay)

    left_glow = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    lg = ImageDraw.Draw(left_glow)
    lg.ellipse((-320, 10, 1120, 1180), fill=(16, 26, 62, 132))
    lg.ellipse((1140, -120, 2140, 620), fill=(10, 82, 124, 78))
    lg.polygon([(0, 760), (640, 340), (1580, 820), (1080, 1080), (0, 1080)], fill=(18, 17, 56, 54))
    left_glow = left_glow.filter(ImageFilter.GaussianBlur(92))
    img = Image.alpha_composite(img, left_glow)

    atmosphere = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    ad = ImageDraw.Draw(atmosphere)
    ad.ellipse((1000, -120, 2180, 720), fill=(40, 82, 255, 42))
    ad.ellipse((1300, 420, 2200, 1160), fill=(98, 69, 255, 34))
    ad.rectangle((0, 0, W, 170), fill=(0, 0, 0, 24))
    ad.rectangle((0, H - 180, W, H), fill=(0, 0, 0, 62))
    atmosphere = atmosphere.filter(ImageFilter.GaussianBlur(76))
    img = Image.alpha_composite(img, atmosphere)

    mesh = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    md = ImageDraw.Draw(mesh, "RGBA")
    for x in range(0, W, 120):
        md.line((x, 0, x, H), fill=(255, 255, 255, 8), width=1)
    for y in range(0, H, 108):
        md.line((0, y, W, y), fill=(255, 255, 255, 6), width=1)
    for x in range(120, W, 240):
        for y in range(96, H, 216):
            md.ellipse((x - 2, y - 2, x + 2, y + 2), fill=(88, 247, 205, 20))
    mesh = mesh.filter(ImageFilter.GaussianBlur(0.4))
    img = Image.alpha_composite(img, mesh)

    draw_cinematic_lights(img)
    return img


def draw_cinematic_lights(canvas: Image.Image, *, intensity: float = 1.0) -> None:
    layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer, "RGBA")
    alpha = int(54 * intensity)
    d.ellipse((1080, 0, 2200, 760), fill=(54, 145, 255, alpha))
    d.ellipse((1180, 360, 2140, 1120), fill=(120, 77, 255, int(30 * intensity)))
    d.ellipse((-260, 680, 880, 1320), fill=(29, 230, 204, int(20 * intensity)))
    d.line((112, 156, 820, 82), fill=(83, 223, 255, int(44 * intensity)), width=4)
    d.line((1080, 984, 1850, 804), fill=(143, 92, 255, int(38 * intensity)), width=4)
    layer = layer.filter(ImageFilter.GaussianBlur(38))
    canvas.alpha_composite(layer)


def gloss_frame(draw: ImageDraw.ImageDraw, margin: int = 54) -> None:
    draw.rounded_rectangle(
        (margin, margin, W - margin, H - margin),
        radius=24,
        outline=(255, 255, 255, 20),
        width=2,
    )
    draw.line((margin + 34, margin, margin + 420, margin), fill=(83, 223, 255, 118), width=4)
    draw.line((W - margin - 460, H - margin, W - margin - 34, H - margin), fill=(143, 92, 255, 110), width=4)
    draw.line((W - margin, margin + 90, W - margin, margin + 320), fill=(88, 247, 205, 70), width=3)


def light_chip(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], text: str, accent=CYAN) -> None:
    x1, y1, x2, y2 = box
    draw.rounded_rectangle((x1, y1, x2, y2), radius=18, fill=(8, 14, 30, 208), outline=accent, width=2)
    draw.rounded_rectangle((x1 + 14, y1 + 12, x1 + 28, y1 + 26), radius=7, fill=accent)
    draw.line((x1 + 48, y1 + 18, x2 - 18, y1 + 18), fill=accent, width=2)
    draw.text((x1 + 48, y1 + 26), text.upper(), font=MONO_15, fill=WHITE)


def paste_logo(canvas: Image.Image, xy: tuple[int, int] = (92, 70), width: int = 214) -> None:
    path = LOGO if LOGO.exists() else MARK
    if not path.exists():
        return
    logo = Image.open(path).convert("RGBA")
    ratio = width / logo.width
    logo = logo.resize((width, int(logo.height * ratio)), Image.Resampling.LANCZOS)
    canvas.alpha_composite(logo, xy)


def cover_crop(path: Path, size: tuple[int, int]) -> Image.Image | None:
    if not path.exists():
        return None
    img = Image.open(path).convert("RGB")
    target_w, target_h = size
    ratio = max(target_w / img.width, target_h / img.height)
    img = img.resize((int(img.width * ratio), int(img.height * ratio)), Image.Resampling.LANCZOS)
    left = max(0, (img.width - target_w) // 2)
    top = max(0, (img.height - target_h) // 2)
    return img.crop((left, top, left + target_w, top + target_h))


def draw_city_horizon(
    canvas: Image.Image,
    *,
    left: int = 0,
    right: int = W,
    base_y: int = H - 38,
    opacity: int = 44,
) -> None:
    layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer, "RGBA")
    widths = [34, 42, 26, 58, 38, 46, 32, 66, 28, 52, 36]
    heights = [84, 132, 98, 178, 122, 156, 110, 214, 92, 164, 118]
    x = left - 18
    index = 0
    while x < right + 40:
        width = widths[index % len(widths)]
        height = heights[index % len(heights)]
        top = base_y - height
        fill = (10, 16, 30, opacity + (index % 4) * 6)
        d.rounded_rectangle((x, top, x + width, base_y), radius=6, fill=fill)
        if index % 2 == 0:
            d.line((x + 6, top + 18, x + width - 6, top + 18), fill=(99, 199, 255, 18), width=2)
        if index % 3 == 0:
            d.line((x + 8, top + 34, x + width - 8, top + 34), fill=(143, 92, 255, 14), width=2)
        if width > 40:
            d.rectangle((x + width // 2 - 2, top - 18, x + width // 2 + 2, top), fill=(99, 199, 255, 42))
            d.ellipse((x + width // 2 - 6, top - 24, x + width // 2 + 6, top - 12), fill=(99, 199, 255, 60))
        x += width - 5
        index += 1

    d.rectangle((left, base_y, right, base_y + 5), fill=(247, 249, 255, 18))
    haze = layer.filter(ImageFilter.GaussianBlur(1.3))
    canvas.alpha_composite(haze)


def portrait_card(canvas: Image.Image, box: tuple[int, int, int, int], path: Path) -> None:
    x1, y1, x2, y2 = box
    layer = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    d = ImageDraw.Draw(layer)
    d.rounded_rectangle(box, radius=34, fill=(8, 13, 25, 228), outline=(85, 116, 175), width=3)
    photo = cover_crop(path, (x2 - x1 - 28, y2 - y1 - 28))
    if photo is None:
        return
    photo = photo.convert("RGBA")
    mask = Image.new("L", photo.size, 0)
    ImageDraw.Draw(mask).rounded_rectangle((0, 0, photo.size[0], photo.size[1]), radius=26, fill=255)
    photo.putalpha(mask)
    glow = Image.new("RGBA", photo.size, (0, 0, 0, 0))
    gd = ImageDraw.Draw(glow)
    gd.rounded_rectangle((0, 0, photo.size[0], photo.size[1]), radius=26, outline=(255, 255, 255, 34), width=2)
    photo.alpha_composite(glow)
    canvas.alpha_composite(layer)
    canvas.alpha_composite(photo, (x1 + 14, y1 + 14))


def draw_text(draw: ImageDraw.ImageDraw, xy: tuple[int, int], text: str, fnt, fill=WHITE, spacing=8) -> None:
    draw.multiline_text(xy, text, font=fnt, fill=fill, spacing=spacing)


def wrap_text(text: str, fnt, max_width: int) -> list[str]:
    words = text.split()
    lines: list[str] = []
    cur = ""
    scratch = Image.new("RGB", (10, 10))
    d = ImageDraw.Draw(scratch)
    for word in words:
        trial = f"{cur} {word}".strip()
        if d.textbbox((0, 0), trial, font=fnt)[2] <= max_width:
            cur = trial
        else:
            if cur:
                lines.append(cur)
            cur = word
    if cur:
        lines.append(cur)
    return lines


def paragraph(draw: ImageDraw.ImageDraw, text: str, x: int, y: int, max_width: int, fnt, fill=SOFT, line_gap: int = 12) -> int:
    lines = wrap_text(text, fnt, max_width)
    line_h = dheight(fnt) + line_gap
    for i, line_text in enumerate(lines):
        draw.text((x, y + i * line_h), line_text, font=fnt, fill=fill)
    return y + len(lines) * line_h


def dheight(fnt) -> int:
    box = fnt.getbbox("Ag")
    return box[3] - box[1]


def label(draw: ImageDraw.ImageDraw, text: str, x: int, y: int, fill=CYAN) -> None:
    draw.text((x, y), text.upper(), font=MONO_18, fill=fill)
    draw.rounded_rectangle((x, y + 34, x + 168, y + 40), radius=3, fill=fill)
    draw.ellipse((x + 176, y + 31, x + 188, y + 43), fill=fill)


def footer(draw: ImageDraw.ImageDraw, idx: int) -> None:
    draw.text((1660, 1018), f"0{idx} / FONATPROP", font=MONO_15, fill=(126, 136, 154))


def rounded_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], fill=CARD, outline=(57, 83, 130), width=2, radius=28) -> None:
    shadow = (box[0], box[1] + 12, box[2], box[3] + 12)
    draw.rounded_rectangle(shadow, radius=radius, fill=(0, 0, 0, 34))
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)
    draw.line((box[0] + 28, box[1] + 18, box[2] - 28, box[1] + 18), fill=(255, 255, 255, 20), width=2)


def video_placeholder(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    tag: str,
    title_text: str,
    accent,
    subtitle: str = "Embedded product video. Click to play in PowerPoint.",
) -> None:
    rounded_card(draw, box, fill=(7, 11, 24, 238), outline=accent, width=3, radius=34)
    x1, y1, x2, y2 = box
    draw.text((x1 + 28, y1 + 24), tag.upper(), font=MONO_15, fill=accent)
    draw.text((x1 + 28, y1 + 58), title_text, font=SANS_B_28, fill=WHITE)
    draw.text((x1 + 28, y2 - 46), subtitle, font=SANS_18, fill=SOFT)
    inner = (x1 + 28, y1 + 110, x2 - 28, y2 - 66)
    draw.rounded_rectangle(inner, radius=24, fill=(13, 21, 40, 192), outline=(72, 100, 156), width=2)
    draw.line((inner[0] + 24, inner[1] + 18, inner[2] - 24, inner[1] + 18), fill=accent, width=2)
    cx = (inner[0] + inner[2]) // 2
    cy = (inner[1] + inner[3]) // 2
    draw.ellipse((cx - 44, cy - 44, cx + 44, cy + 44), fill=(248, 250, 252, 230))
    draw.polygon([(cx - 10, cy - 16), (cx - 10, cy + 16), (cx + 18, cy)], fill=INK)


def copy_block(
    draw: ImageDraw.ImageDraw,
    box: tuple[int, int, int, int],
    text: str,
    *,
    body_font=SANS_24,
    fill=(8, 13, 25, 182),
    outline=(56, 82, 126),
    text_fill=SOFT,
    radius: int = 28,
    pad_x: int = 30,
    pad_y: int = 28,
    line_gap: int = 10,
) -> None:
    rounded_card(draw, box, fill=fill, outline=outline, radius=radius)
    x1, y1, x2, _ = box
    paragraph(draw, text, x1 + pad_x, y1 + pad_y, x2 - x1 - pad_x * 2, body_font, text_fill, line_gap)


def metric(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], value: str, cap: str) -> None:
    rounded_card(draw, box, fill=(8, 12, 26, 198), outline=(69, 104, 188), radius=20)
    x1, y1, x2, _ = box
    draw.text((x1 + 24, y1 + 22), value, font=SANS_B_34, fill=WHITE)
    draw.text((x1 + 24, y1 + 84), cap.upper(), font=MONO_15, fill=MUTED)


def bullet_panel(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], title_text: str, bullets: list[str]) -> None:
    rounded_card(draw, box, fill=(8, 13, 25, 228), outline=(74, 109, 171), radius=38)
    x1, y1, x2, _ = box
    draw.text((x1 + 44, y1 + 42), title_text.upper(), font=MONO_18, fill=BLUE)
    y = y1 + 108
    for i, bullet in enumerate(bullets):
        draw.ellipse((x1 + 44, y + 10, x1 + 56, y + 22), fill=BLUE if i == 0 else CYAN)
        y = paragraph(draw, bullet, x1 + 84, y, x2 - x1 - 132, SANS_24, WHITE if i == 0 else SOFT, 8) + 20


def product_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], tag: str, title_text: str, copy: str, accent) -> None:
    rounded_card(draw, box, fill=CARD2, outline=accent, radius=24, width=3)
    x1, y1, x2, _ = box
    draw.text((x1 + 34, y1 + 30), tag.upper(), font=MONO_15, fill=accent)
    draw.text((x1 + 34, y1 + 78), title_text, font=SERIF_56, fill=WHITE)
    paragraph(draw, copy, x1 + 34, y1 + 156, x2 - x1 - 68, SANS_22, SOFT, 8)


def step_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], number: str, title_text: str, copy: str) -> None:
    rounded_card(draw, box, fill=(8, 13, 25, 214), outline=(67, 97, 154), radius=22)
    x1, y1, x2, _ = box
    draw.text((x1 + 28, y1 + 22), number, font=MONO_18, fill=BLUE)
    draw.text((x1 + 84, y1 + 14), title_text, font=SERIF_56, fill=WHITE)
    paragraph(draw, copy, x1 + 28, y1 + 78, x2 - x1 - 56, SANS_20, SOFT, 8)


def pricing_card(draw: ImageDraw.ImageDraw, box: tuple[int, int, int, int], name: str, aed: str, usd: str, accent, featured=False) -> None:
    rounded_card(draw, box, fill=(8, 13, 25, 235), outline=accent, width=3 if featured else 2, radius=26)
    x1, y1, x2, _ = box
    if featured:
        draw.rounded_rectangle((x1 + 28, y1 + 24, x1 + 174, y1 + 58), radius=8, fill=accent)
        draw.text((x1 + 58, y1 + 32), "POPULAR", font=MONO_15, fill=INK)
        name_y = y1 + 82
    else:
        name_y = y1 + 45
    draw.text((x1 + 28, name_y), name, font=SANS_B_22, fill=SOFT)
    price_font = SERIF_56 if (x2 - x1) >= 300 and len(aed) <= 9 else SERIF_66 if len(aed) <= 8 else SANS_B_34
    if (x2 - x1) < 280:
        price_font = SANS_B_34
    price_box = draw.textbbox((0, 0), aed, font=price_font)
    price_w = price_box[2] - price_box[0]
    draw.text((x1 + ((x2 - x1) - price_w) / 2, y1 + 128), aed, font=price_font, fill=WHITE)
    usd_box = draw.textbbox((0, 0), usd, font=SANS_20)
    usd_w = usd_box[2] - usd_box[0]
    draw.text((x1 + ((x2 - x1) - usd_w) / 2, y1 + 206), usd, font=SANS_20, fill=MUTED)


def slide_1() -> Image.Image:
    img = premium_bg(BG_IMAGES[0], 80)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img, (92, 72), 230)
    label(draw, "Dubai brokerages", 96, 210)
    draw.text((96, 282), "AI valuation.", font=SERIF_88, fill=WHITE)
    draw.text((96, 382), "Qualified leads.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 512, 1040, 694),
        "Private AI valuations. Public lead capture. Backed by 234K+ real Dubai transactions.",
        body_font=SANS_28,
        fill=(8, 13, 25, 168),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=30,
        line_gap=12,
    )
    metric(draw, (100, 758, 368, 904), "234K+", "real transactions")
    metric(draw, (408, 758, 676, 904), "AI", "valuation model")
    metric(draw, (716, 758, 984, 904), "Widget", "lead capture")
    light_chip(draw, (1188, 806, 1628, 864), "Official evidence + AI workflow", CYAN)
    footer(draw, 1)
    return img


def slide_2() -> Image.Image:
    img = premium_bg(BG_IMAGES[1], 74)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Founders", 96, 184)
    draw.text((96, 250), "Who builds FonatProp.", font=SANS_B_74, fill=WHITE)
    draw.text((96, 348), "Real operators.", font=SANS_52, fill=SOFT)
    paragraph(
        draw,
        "A young team building AI real estate tools, live in Dubai and expanding into France.",
        96,
        446,
        760,
        SANS_28,
        SOFT,
        10,
    )

    rounded_card(draw, (1116, 214, 1812, 444), fill=(8, 13, 25, 214), outline=(73, 108, 171), radius=34, width=3)
    draw.text((1160, 254), "WHY TRUST MATTERS", font=MONO_18, fill=CYAN)
    paragraph(
        draw,
        "Agencies should know who is behind the platform they buy. FonatProp is being built by real operators, with product live in Dubai and France already in build.",
        1160,
        312,
        596,
        SANS_24,
        SOFT,
        9,
    )

    portrait_card(img, (96, 590, 362, 964), FOUNDER_PHOTO)

    rounded_card(draw, (392, 590, 1090, 964), fill=(8, 13, 25, 228), outline=(73, 108, 171), radius=34, width=3)
    draw.text((440, 636), "FOUNDER PROFILE", font=MONO_18, fill=CYAN)
    draw.text((440, 692), "Francisco Natale Oertly / Founder", font=SANS_B_34, fill=WHITE)
    founder_lines = [
        "22 years old / Professional footballer",
        "Sagrado Corazon / High school",
        "UAI University / 4th year",
        "Building AI real estate tools with a young operating team",
    ]
    y = 766
    for item in founder_lines:
        draw.ellipse((444, y + 8, 456, y + 20), fill=CYAN)
        draw.text((480, y), item, font=SANS_22, fill=SOFT)
        y += 42

    milestone_cards = [
        ("01 / Team", "Young operators", "Built by a small team focused on real estate execution, not generic agency work.", CYAN),
        ("02 / Dubai", "Live product", "Dubai already combines AI valuation, lead capture and a broker workflow.", GOLD),
        ("03 / France", "Second market", "France is already in build as the next FonatProp market surface.", VIOLET),
    ]
    x = 1120
    for label_text, title_text, copy, accent in milestone_cards:
        rounded_card(draw, (x, 520, x + 212, 964), fill=(8, 13, 25, 228), outline=accent, radius=28, width=3)
        draw.text((x + 22, 554), label_text.upper(), font=MONO_15, fill=accent)
        draw.text((x + 22, 608), title_text, font=SANS_B_22, fill=WHITE)
        paragraph(draw, copy, x + 22, 654, 168, SANS_20, SOFT, 6)
        x += 232
    footer(draw, 2)
    return img


def slide_3() -> Image.Image:
    img = premium_bg(BG_IMAGES[1], 72)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Pain", 96, 214)
    draw.text((96, 292), "Slow valuations.", font=SERIF_88, fill=WHITE)
    draw.text((96, 392), "Lost website leads.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 532, 988, 712),
        "Many brokerages still rely on opinion, manual comps and slow follow-up. Website traffic leaves without converting.",
        body_font=SANS_28,
        fill=(8, 13, 25, 166),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=28,
        line_gap=12,
    )
    bullet_panel(
        draw,
        (1190, 226, 1800, 762),
        "Brokerage pain",
        [
            "Slow pricing decisions",
            "Subjective broker opinions",
            "Anonymous website traffic",
            "Late follow-up",
        ],
    )
    light_chip(draw, (1190, 810, 1704, 868), "Problem: value + conversion speed", GOLD)
    footer(draw, 3)
    return img


def slide_4() -> Image.Image:
    img = premium_bg(BG_IMAGES[2], 76)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Private AI", 96, 214)
    draw.text((96, 292), "Fast valuations.", font=SERIF_88, fill=WHITE)
    draw.text((96, 392), "Real evidence.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 540, 860, 724),
        "Address in. AI range out. Guided by real Dubai transaction evidence and broker review logic.",
        body_font=SANS_28,
        fill=(8, 13, 25, 168),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=30,
        line_gap=12,
    )
    video_placeholder(
        draw,
        (930, 210, 1812, 886),
        "Private AI workflow",
        "AI valuation demo",
        BLUE,
        subtitle="Click to play in PowerPoint.",
    )
    footer(draw, 4)
    return img


def slide_5() -> Image.Image:
    img = premium_bg(BG_IMAGES[3], 76)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Public AI", 96, 214)
    draw.text((96, 292), "Website widget.", font=SERIF_88, fill=WHITE)
    draw.text((96, 392), "Lead first.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 540, 860, 724),
        "Capture contact details first. Show a broad AI range second. Keep the final follow-up inside the agency.",
        body_font=SANS_28,
        fill=(8, 13, 25, 168),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=30,
        line_gap=12,
    )
    video_placeholder(
        draw,
        (930, 210, 1812, 886),
        "Public AI workflow",
        "Widget demo",
        GOLD,
        subtitle="Click to play in PowerPoint.",
    )
    footer(draw, 5)
    return img


def slide_6() -> Image.Image:
    img = premium_bg(BG_IMAGES[4], 70)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img, (92, 64), 206)
    label(draw, "Plans", 96, 156, GOLD)
    draw.text((96, 232), "Simple pricing.", font=SERIF_78, fill=WHITE)
    draw.text((96, 320), "Easy to buy.", font=SERIF_I_58, fill=SOFT)
    top_cards = [
        ("Widget", "Fixed monthly fee", "Lead capture"),
        ("Valuation", "1 search = 1 token", "Private AI usage"),
        ("Top-up", "Extra credits on request", "Buy more or upgrade"),
    ]
    x = 930
    for title_text, detail, sub in top_cards:
        rounded_card(draw, (x, 228, x + 266, 340), fill=(11, 18, 33, 226), outline=(62, 94, 153), radius=22)
        draw.text((x + 18, 250), title_text.upper(), font=MONO_15, fill=CYAN)
        draw.text((x + 18, 284), detail, font=SANS_B_18, fill=WHITE)
        draw.text((x + 18, 314), sub, font=SANS_18, fill=SOFT)
        x += 286

    rounded_card(draw, (96, 408, 520, 932), fill=(8, 13, 25, 236), outline=GOLD, width=4, radius=38)
    draw.rounded_rectangle((128, 446, 266, 482), radius=10, fill=GOLD)
    draw.text((154, 468), "WIDGET", font=MONO_15, fill=INK)
    draw.text((128, 532), "Lead capture", font=SANS_B_34, fill=WHITE)
    draw.text((128, 598), "AED 1,099", font=SERIF_66, fill=WHITE)
    draw.text((132, 668), "$299 / month", font=SANS_24, fill=SOFT)
    paragraph(draw, "One website. One widget. One monthly fee.", 128, 734, 330, SANS_20, SOFT, 6)
    rounded_card(draw, (128, 826, 488, 892), fill=(11, 18, 33, 236), outline=(88, 115, 171), radius=18)
    draw.text((154, 848), "Single-site license", font=SANS_B_18, fill=WHITE)
    draw.text((358, 848), "Live", font=MONO_15, fill=GOLD)

    plans_box = (580, 392, 1812, 940)
    rounded_card(draw, plans_box, fill=(8, 13, 25, 238), outline=(62, 94, 153), width=3, radius=38)
    draw.text((628, 432), "AI valuation plans", font=SANS_B_34, fill=WHITE)
    draw.text((628, 478), "Monthly AI access for brokerage teams.", font=SANS_22, fill=SOFT)
    draw.text((1382, 444), "10 valuations / AED 1,465", font=SANS_B_22, fill=WHITE)
    draw.text((1412, 476), "Scales up to 200 valuations", font=SANS_18, fill=SOFT)

    rows = [
        ("10 valuations / month", "AED 1,465", "$399 / mo", False),
        ("20 valuations / month", "AED 2,199", "$599 / mo", True),
        ("50 valuations / month", "AED 4,400", "$1,199 / mo", False),
        ("100 valuations / month", "AED 7,345", "$2,000 / mo", False),
        ("200 valuations / month", "AED 12,855", "$3,500 / mo", False),
    ]
    y = 534
    for title_text, aed, usd, featured in rows:
        fill = (14, 22, 38, 232) if not featured else (16, 24, 40, 244)
        outline = GOLD if featured else (50, 74, 122)
        draw.rounded_rectangle((628, y, 1768, y + 60), radius=22, fill=fill, outline=outline, width=3 if featured else 2)
        if featured:
            draw.rounded_rectangle((648, y + 15, 744, y + 43), radius=10, fill=GOLD)
            draw.text((670, y + 22), "POPULAR", font=MONO_15, fill=INK)
            title_x = 772
        else:
            title_x = 660
        draw.text((title_x, y + 16), title_text, font=SANS_B_22, fill=WHITE)
        aed_box = draw.textbbox((0, 0), aed, font=SANS_B_28)
        aed_w = aed_box[2] - aed_box[0]
        usd_box = draw.textbbox((0, 0), usd, font=SANS_20)
        usd_w = usd_box[2] - usd_box[0]
        draw.text((1464 - aed_w / 2, y + 12), aed, font=SANS_B_28, fill=WHITE)
        draw.text((1662 - usd_w / 2, y + 20), usd, font=SANS_20, fill=SOFT)
        y += 72

    draw.text((628, 900), "Extra credits, custom bundles and enterprise setup on request.", font=SANS_18, fill=SOFT)
    footer(draw, 6)
    return img


def slide_7() -> Image.Image:
    img = premium_bg(BG_IMAGES[5], 74)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Commercial ROI", 96, 190, GOLD)
    draw.text((96, 268), "One mandate", font=SERIF_88, fill=WHITE)
    draw.text((96, 368), "can pay back.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 520, 870, 704),
        "If the widget turns website traffic into one serious seller conversation, the monthly stack becomes easy to justify.",
        body_font=SANS_28,
        fill=(8, 13, 25, 168),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=28,
        line_gap=12,
    )

    rounded_card(draw, (990, 192, 1812, 872), fill=(8, 13, 25, 236), outline=GOLD, width=4, radius=42)
    draw.text((1044, 242), "ROI EXAMPLE", font=MONO_18, fill=GOLD)
    draw.text((1044, 300), "AED 3,298 / mo", font=SERIF_66, fill=WHITE)
    draw.text((1048, 372), "Widget AED 1,099 + 20 valuation tokens AED 2,199", font=SANS_20, fill=SOFT)

    roi_rows = [
        ("1,000 website visits", "30 captured leads", "3.0% conversion"),
        ("30 leads", "9 broker conversations", "30% contacted"),
        ("9 conversations", "1 mandate", "seller won"),
        ("AED 1.5M sale", "AED 30K commission", "2% agency fee"),
    ]
    y = 448
    for left, right, tag in roi_rows:
        draw.rounded_rectangle((1044, y, 1754, y + 74), radius=22, fill=(11, 18, 33, 232), outline=(62, 94, 153), width=2)
        draw.text((1072, y + 22), left, font=SANS_B_22, fill=WHITE)
        draw.text((1374, y + 22), right, font=SANS_B_22, fill=CYAN if "AED 30K" not in right else GOLD)
        draw.text((1640, y + 27), tag.upper(), font=MONO_15, fill=MUTED)
        y += 92

    light_chip(draw, (1044, 800, 1608, 862), "Example scenario, not a guarantee", CYAN)
    metric(draw, (96, 760, 386, 906), "9x", "monthly payback")
    metric(draw, (426, 760, 716, 906), "1", "mandate target")
    footer(draw, 7)
    return img


def slide_8() -> Image.Image:
    img = premium_bg(BG_IMAGES[5], 74)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img)
    label(draw, "Next step", 96, 214)
    draw.text((96, 292), "Open the demo.", font=SERIF_88, fill=WHITE)
    draw.text((96, 392), "See it live.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 528, 1038, 702),
        "Test the AI valuation flow. Watch the widget capture the lead. Review the product live.",
        body_font=SANS_28,
        fill=(8, 13, 25, 162),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=28,
        line_gap=12,
    )
    rounded_card(draw, (96, 748, 1096, 972), fill=(8, 13, 25, 222), outline=(66, 98, 154), radius=30)
    draw.text((132, 786), "LIVE BROKER DEMO", font=MONO_18, fill=CYAN)
    draw.rounded_rectangle((132, 824, 676, 914), radius=24, fill=WHITE)
    draw.text((190, 856), "OPEN LIVE BROKER DEMO", font=SANS_B_22, fill=INK)
    draw.text((132, 934), DEMO_URL, font=SANS_B_18, fill=CYAN)
    if QR.exists():
        rounded_card(draw, (1388, 674, 1718, 972), fill=(8, 13, 25, 232), outline=BLUE, radius=30)
        qr = Image.open(QR).convert("RGBA").resize((214, 214), Image.Resampling.LANCZOS)
        img.alpha_composite(qr, (1446, 712))
        draw.text((1492, 942), "SCAN DEMO", font=MONO_15, fill=CYAN)
    footer(draw, 8)
    return img


def slide_9() -> Image.Image:
    img = premium_bg(BG_IMAGES[6], 72)
    draw = ImageDraw.Draw(img)
    gloss_frame(draw)
    paste_logo(img, (92, 72), 230)
    label(draw, "Contact", 96, 214, GOLD)
    draw.text((96, 292), "Let's talk.", font=SERIF_88, fill=WHITE)
    draw.text((96, 392), "FonatProp.", font=SERIF_I_76, fill=SOFT)
    copy_block(
        draw,
        (96, 526, 982, 704),
        "For brokerages, private demos and partnership conversations.",
        body_font=SANS_32,
        fill=(8, 13, 25, 164),
        outline=(61, 92, 144),
        pad_x=32,
        pad_y=32,
        line_gap=12,
    )

    rounded_card(draw, (1096, 214, 1812, 892), fill=(8, 13, 25, 234), outline=GOLD, width=4, radius=42)
    draw.text((1148, 262), "DIRECT CONTACT", font=MONO_18, fill=GOLD)
    draw.line((1148, 316, 1710, 316), fill=(235, 196, 105, 90), width=2)

    contact_rows = [
        ("WHATSAPP", CONTACT_PHONE_DISPLAY, "Argentina mobile code included"),
        ("EMAIL", CONTACT_EMAIL, "Commercial contact"),
        ("LIVE DEMO", DEMO_URL, "Broker demo experience"),
    ]
    y = 374
    for tag, value, sub in contact_rows:
        rounded_card(draw, (1148, y, 1748, y + 122), fill=(11, 18, 33, 228), outline=(64, 93, 150), radius=24)
        draw.text((1184, y + 26), tag, font=MONO_15, fill=CYAN if tag != "WHATSAPP" else GOLD)
        draw.text((1184, y + 58), value, font=SANS_B_28 if tag != "LIVE DEMO" else SANS_B_22, fill=WHITE)
        draw.text((1184, y + 92), sub, font=SANS_18, fill=SOFT)
        y += 150

    light_chip(draw, (96, 780, 634, 842), "AI valuation + lead capture for real estate", CYAN)
    light_chip(draw, (96, 862, 514, 924), "Dubai live. France commercial beta.", GOLD)
    footer(draw, 9)
    return img


def add_transitions(pptx_path: Path) -> None:
    ns = {"p": "http://schemas.openxmlformats.org/presentationml/2006/main"}
    ET.register_namespace("p", ns["p"])
    tmp = pptx_path.with_suffix(".tmp.pptx")
    with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                root = ET.fromstring(data)
                for old in root.findall("p:transition", ns):
                    root.remove(old)
                transition = ET.Element(f"{{{ns['p']}}}transition", {"spd": "med", "advClick": "1"})
                ET.SubElement(transition, f"{{{ns['p']}}}fade")
                root.insert(1, transition)
                data = ET.tostring(root, encoding="utf-8", xml_declaration=True)
            zout.writestr(item, data)
    tmp.replace(pptx_path)


def ensure_video_posters() -> None:
    try:
        import cv2  # type: ignore
    except Exception:
        return

    DECK_MEDIA.mkdir(parents=True, exist_ok=True)
    targets = [
        (VALUATION_VIDEO, VALUATION_POSTER, 0.60),
        (WIDGET_VIDEO, WIDGET_POSTER, 0.08),
    ]
    for movie_path, poster_path, ratio in targets:
        if not movie_path.exists():
            continue
        cap = cv2.VideoCapture(str(movie_path))
        if not cap.isOpened():
            continue
        frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        index = max(0, min(frames - 1, int(frames * ratio)))
        cap.set(cv2.CAP_PROP_POS_FRAMES, index)
        ok, frame = cap.read()
        if ok:
            cv2.imwrite(str(poster_path), frame)
        cap.release()


def make_deck(slide_paths: list[Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]
    video_embeds = {
        4: {
            "movie": VALUATION_VIDEO,
            "poster": VALUATION_POSTER,
            "left": Inches(6.66),
            "top": Inches(2.26),
            "width": Inches(5.48),
            "height": Inches(3.08),
        },
        5: {
            "movie": WIDGET_VIDEO,
            "poster": WIDGET_POSTER,
            "left": Inches(6.66),
            "top": Inches(2.26),
            "width": Inches(5.48),
            "height": Inches(3.08),
        },
    }
    for idx, path in enumerate(slide_paths, 1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(path), 0, 0, width=SLIDE_W, height=SLIDE_H)
        if idx in video_embeds:
            spec = video_embeds[idx]
            movie_path = spec["movie"]
            poster_path = spec["poster"]
            if movie_path.exists() and poster_path.exists():
                slide.shapes.add_movie(
                    str(movie_path),
                    spec["left"],
                    spec["top"],
                    spec["width"],
                    spec["height"],
                    poster_frame_image=str(poster_path),
                    mime_type="video/mp4",
                )
        if idx == 8:
            link = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.92), Inches(5.72), Inches(3.78), Inches(0.74))
            link.fill.solid()
            link.fill.fore_color.rgb = RGBColor(248, 250, 252)
            link.line.color.rgb = RGBColor(248, 250, 252)
            link.text_frame.clear()
            paragraph_shape = link.text_frame.paragraphs[0]
            paragraph_shape.alignment = PP_ALIGN.CENTER
            run = paragraph_shape.add_run()
            run.text = "OPEN LIVE BROKER DEMO"
            run.font.bold = True
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(3, 5, 10)
            link.click_action.hyperlink.address = DEMO_URL
        if idx == 9:
            hotspots = [
                (Inches(7.78), Inches(2.6), Inches(4.17), Inches(0.86), CONTACT_WHATSAPP_URL),
                (Inches(7.78), Inches(3.64), Inches(4.17), Inches(0.86), CONTACT_EMAIL_URL),
                (Inches(7.78), Inches(4.68), Inches(4.17), Inches(0.86), DEMO_URL),
            ]
            for left, top, width, height, url in hotspots:
                link = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, width, height)
                link.fill.background()
                link.line.fill.background()
                link.click_action.hyperlink.address = url
    return prs


def main() -> None:
    ensure_video_posters()
    tmpdir = Path(tempfile.gettempdir()) / "fonatprop_premium_deck"
    if tmpdir.exists():
        shutil.rmtree(tmpdir)
    tmpdir.mkdir(parents=True, exist_ok=True)

    slides = [
        slide_1(),
        slide_2(),
        slide_3(),
        slide_4(),
        slide_5(),
        slide_6(),
        slide_7(),
        slide_8(),
        slide_9(),
    ]
    slide_paths = []
    for idx, slide in enumerate(slides, 1):
        path = tmpdir / f"slide_{idx:02}.png"
        slide.convert("RGB").save(path, quality=96)
        slide_paths.append(path)

    outputs = [OUT_DESKTOP, OUT_DOWNLOADS, OUT_DATA_LAKE]
    for output in outputs:
        output.parent.mkdir(parents=True, exist_ok=True)
        deck = make_deck(slide_paths)
        deck.save(output)
        add_transitions(output)
        print(f"Saved {output}")


if __name__ == "__main__":
    main()
